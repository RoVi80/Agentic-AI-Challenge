"""Inspect a reference .pptx and return its slide layouts as structured data.

Companion tool to brand_render. Where brand_render takes a chosen layout and
fills it, inspect_brand_layouts gives the agent the *menu* to choose from.

For each layout the reference deck offers, returns:
  - index, name, master_name
  - placeholder_count
  - placeholders: list of {idx, type, name}

The agent calls this first whenever a user uploads a reference. It then picks
a layout BY NAME per slide and passes that name to brand_render — replacing
the keyword heuristic the renderer used to do internally.

Run locally:
    python tools/inspect_brand_layouts.py <reference.pptx>
"""
from ibm_watsonx_orchestrate.agent_builder.tools import tool, WXOFile

import argparse
import io
import json
from pathlib import Path

from pptx import Presentation


def _placeholder_type_name(ph):
    """Return the placeholder's type as the standard PowerPoint name
    ('TITLE', 'BODY', 'PICTURE', 'OBJECT', etc.).

    These names come from python-pptx's PP_PLACEHOLDER enum — itself a direct
    mapping of the OOXML / PresentationML standard. We're not labeling anything
    ourselves; we're surfacing the format's own vocabulary for the LLM to read."""
    try:
        ph_type = ph.placeholder_format.type
    except Exception:
        return "UNKNOWN"
    if ph_type is None:
        return "UNKNOWN"
    return getattr(ph_type, "name", "UNKNOWN")


def _placeholder_info(ph):
    try:
        idx = ph.placeholder_format.idx
    except Exception:
        idx = None
    return {
        "idx": idx,
        "type": _placeholder_type_name(ph),
        "name": ph.name or "",
    }


def _master_name(layout):
    try:
        return layout.slide_master.name or ""
    except Exception:
        return ""


# Placeholder-type categories used to derive body_placeholder_count.
# Kept in lockstep with brand_render._TITLE_TYPES / _SKIP_FILL_TYPES so the
# count we surface is exactly what the renderer can actually fill.
_TITLE_TYPES = {"TITLE", "CENTER_TITLE"}
_SKIP_FILL_TYPES = {"DATE", "SLIDE_NUMBER", "FOOTER", "HEADER", "PICTURE", "SLIDE_IMAGE"}


def _body_placeholder_count(placeholders):
    """Count placeholders that the renderer treats as content slots — i.e.
    BODY, OBJECT, SUBTITLE. Excludes the title slot, master-inherited slots
    (date/footer/slide_number/header), and image slots (picture/slide_image).

    This is the number the agent should compare against its own slide content:
    a 3-column slide needs body_placeholder_count >= 3, otherwise the renderer
    will silently drop the extra columns."""
    count = 0
    for ph in placeholders:
        t = ph.get("type", "")
        if t in _TITLE_TYPES or t in _SKIP_FILL_TYPES:
            continue
        count += 1
    return count


def inspect_layouts_from_bytes(reference_bytes):
    """Return {"layouts": [...]} for the given .pptx bytes.

    Mirrors brand_render's view of the deck: iterates prs.slide_layouts only
    (the first master's layouts). If the deck has multiple slide masters,
    layouts from secondary masters are not exposed yet — same scope as the
    renderer, kept in lockstep on purpose.

    Each layout entry has an `is_first` boolean — true only for the layout at
    index 0. Brand templates almost universally put their default opener
    layout first, so this is a strong prior for "use this for the title slide"
    without any keyword matching."""
    prs = Presentation(io.BytesIO(reference_bytes))
    layouts_out = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = [_placeholder_info(ph) for ph in layout.placeholders]
        layouts_out.append({
            "index": i,
            "name": layout.name or "",
            "master_name": _master_name(layout),
            "is_first": i == 0,
            "placeholder_count": len(placeholders),
            "body_placeholder_count": _body_placeholder_count(placeholders),
            "placeholders": placeholders,
        })
    return {"layouts": layouts_out}


@tool(
    name="inspect_brand_layouts",
    description=(
        "Inspects a reference .pptx and returns its slide layouts so the agent "
        "can pick the right one per slide. For each layout, returns: name, "
        "master_name, is_first (true only for layout index 0 — usually the "
        "branded opener), placeholder_count, body_placeholder_count (count of "
        "fillable content slots — BODY/OBJECT/SUBTITLE — excluding the title "
        "and master-inherited slots like footer/date/picture), and the full "
        "placeholders list. The agent MUST verify body_placeholder_count >= the "
        "number of body sections its slide contains (1 for a single bullet "
        "block, 2 for two-column, 3 for three-column, etc.) — extras are "
        "silently dropped by the renderer. Call this BEFORE brand_render "
        "whenever a reference is uploaded — the layout names returned here are "
        "exactly what brand_render's per-slide 'layout_name' field expects."
    ),
)
def inspect_brand_layouts(pptx_file: WXOFile) -> dict:
    reference_bytes = WXOFile.get_content(pptx_file)
    return inspect_layouts_from_bytes(reference_bytes)


# ─── CLI for local testing ────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("reference_pptx", help="Path to the brand's reference .pptx")
    args = parser.parse_args()

    ref_bytes = Path(args.reference_pptx).read_bytes()
    result = inspect_layouts_from_bytes(ref_bytes)
    print(json.dumps(result, indent=2, ensure_ascii=False))
