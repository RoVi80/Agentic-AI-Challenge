"""Brand-native PowerPoint renderer.

Instead of drawing slides from scratch with our own layout vocabulary, this tool
opens the user's reference .pptx, deletes its existing slides, and adds NEW slides
using the reference's own slide layouts (master templates). The output inherits
every visual decision the brand already made: logos, colors, fonts, footer
treatments, accent shapes, layout grids — all without us extracting or redrawing
any of it.

Used by the rolando_style_clone_agent. Also runnable as CLI for local testing.

Run locally:
    python tools/brand_render.py <reference.pptx> <spec.json> <output.pptx>

The spec JSON shape (no settings needed — the reference deck IS the spec).

Primary form — pick a layout by exact name from inspect_brand_layouts output:
    {
      "slides": [
        { "layout_name": "Background 1",   "title": "...", "subtitle": "..." },
        { "layout_name": "Titel und Inhalt", "title": "...", "bullets": ["...", "..."] },
        { "layout_name": "Drei Inhalte",   "title": "...",
          "bodies": [ ["bullet a", "bullet b"], "single line", ["x", "y"] ] },
        { "layout_name": "Zitat",          "title": "...", "bullets": ["..."],
          "layout_reason": "quote slide for the testimonial" }
      ]
    }

Legacy/fallback form — used when layout_name is missing or doesn't match any
layout. The renderer's heuristic picks among the brand's layouts using the type:
    {
      "slides": [
        { "type": "title",      "title": "...", "subtitle": "..." },
        { "type": "bullets",    "title": "...", "bullets": ["...", "..."] },
        { "type": "two_column", "title": "...",
          "leftTitle": "...", "leftBullets": ["..."],
          "rightTitle": "...", "rightBullets": ["..."] },
        { "type": "closing",    "title": "...", "bullets": ["..."] }
      ]
    }
"""
from ibm_watsonx_orchestrate.agent_builder.tools import tool, WXOFile

import argparse
import io
import json
import re
import zipfile
from pathlib import Path

from pptx import Presentation


# ─── Layout matching ──────────────────────────────────────────────────────

def _placeholder_count(layout):
    """Number of placeholders the layout exposes."""
    try:
        return len(list(layout.placeholders))
    except Exception:
        return 0


def _name_lower(layout):
    return (getattr(layout, "name", "") or "").lower()


def _is_picture_placeholder(ph):
    """An image/picture placeholder. python-pptx mis-types these in some decks,
    so we also fall back to a name check covering common languages."""
    try:
        t = int(ph.placeholder_format.type)
    except Exception:
        t = 0
    name = (getattr(ph, "name", "") or "").lower()
    return t == 28 or any(tok in name for tok in ("picture", "image", "bild"))


# UNIVERSAL pattern: section/chapter divider layouts contain huge single-glyph
# placeholders ("A", "1", a giant chapter letter) that mangle ordinary text.
# This is the only universally-true layout pattern we lean on, because every
# template family has a "section divider" type, and they're never appropriate
# for "title", "bullets", "two_column", or "closing" slides.
_SECTION_DIVIDER_TOKENS = (
    "section", "chapter", "chapitre", "kapitel", "divider", "kapiteltitel",
)


def _score_layout(layout, slide_type, layout_index):
    """Pragmatic heuristic: placeholder count + a small set of universal English
    keywords (PowerPoint's stock layouts are in English) + one universal rule
    (never pick a section divider for a content slide). Imperfect — the durable
    answer is to let the LLM pick from the layout list."""
    name = _name_lower(layout)
    pcount = _placeholder_count(layout)
    score = 0

    # Universal: section/chapter dividers are never the right pick for any
    # of our four slide types (they exist to introduce a section, not hold content)
    if any(tok in name for tok in _SECTION_DIVIDER_TOKENS):
        score -= 15

    if slide_type == "title":
        if layout_index == 0:
            score += 6
        if "title" in name:
            score += 3
        if pcount <= 6:
            score += 1

    elif slide_type == "bullets":
        if "title and content" in name or "content" in name:
            score += 5
        if pcount in (2, 6):
            score += 3
        if layout_index == 0:
            score -= 4

    elif slide_type == "two_column":
        if "two content" in name or "comparison" in name:
            score += 6
        if pcount in (3, 6, 7):
            score += 4
        if layout_index == 0:
            score -= 4

    elif slide_type == "closing":
        if "content" in name and pcount in (2, 6):
            score += 3
        if layout_index == 0:
            score -= 3

    return score


def _pick_layout(slide_type, layouts):
    """Return the best-matching brand layout for our simple slide type."""
    scored = [
        (_score_layout(layout, slide_type, i), i, layout)
        for i, layout in enumerate(layouts)
    ]
    scored.sort(key=lambda x: (-x[0], x[1]))
    if scored and scored[0][0] > 0:
        return scored[0][2]
    # No good match — return the second layout (often a generic content layout)
    if len(layouts) >= 2:
        return layouts[1]
    return layouts[0] if layouts else None


def _lookup_layout_by_name(name, layouts):
    """Find a layout by exact name (case-insensitive). Returns None if no match
    so the caller can fall back to the heuristic — we deliberately do NOT
    fuzzy-match, because a wrong fuzzy match is worse than a clean fallback."""
    if not name:
        return None
    target = str(name).strip().lower()
    for layout in layouts:
        if (getattr(layout, "name", "") or "").strip().lower() == target:
            return layout
    return None


# ─── Placeholder filling ──────────────────────────────────────────────────

def _set_text(placeholder, text):
    """Set a placeholder's text, preserving the layout's default formatting."""
    if placeholder is None or text is None:
        return
    try:
        placeholder.text = str(text)
    except Exception:
        pass


def _set_bullets(placeholder, items):
    """Fill a placeholder with bullet items. Each item becomes its own paragraph
    so the layout's bullet style applies."""
    if placeholder is None or not items:
        return
    try:
        tf = placeholder.text_frame
        tf.text = str(items[0]) if items else ""
        for item in items[1:]:
            p = tf.add_paragraph()
            p.text = str(item)
    except Exception:
        pass


def _remove_empty_picture_placeholders(slide):
    """After we've filled what we want, remove any leftover image placeholders
    so they don't render as 'drag image here' transparent grids."""
    to_remove = []
    for ph in slide.placeholders:
        if _is_picture_placeholder(ph):
            to_remove.append(ph)
    for ph in to_remove:
        try:
            sp = ph._element
            sp.getparent().remove(sp)
        except Exception:
            pass


def _remove_unfilled_text_placeholders(slide):
    """Strip BODY/OBJECT/SUBTITLE placeholders that we didn't fill. Without
    this, an empty placeholder renders as the layout's default prompt — e.g.
    'Click to add text', or in a French-authored deck 'Espace réservé du
    texte'. These are visual noise from the brand's master, not our content,
    so they should disappear when we don't have anything to put there."""
    to_remove = []
    for ph in slide.placeholders:
        try:
            type_name = ph.placeholder_format.type
            type_name = getattr(type_name, "name", "") or ""
        except Exception:
            type_name = ""
        # Skip title (always filled with our content or intentionally left
        # to inherit), master-inherited slots, and pictures (handled above).
        if type_name in _TITLE_TYPES or type_name in _SKIP_FILL_TYPES:
            continue
        if _is_picture_placeholder(ph):
            continue
        # Empty content slot — we never set text on it, so it would render
        # the master's default prompt.
        try:
            text = (ph.text_frame.text or "").strip()
        except Exception:
            text = ""
        if not text:
            to_remove.append(ph)
    for ph in to_remove:
        try:
            sp = ph._element
            sp.getparent().remove(sp)
        except Exception:
            pass


def _title_and_body_placeholders(slide):
    """Return (title_placeholder, body_placeholder_list) ordered by appearance.
    Title is the first placeholder identified as a title type, or just the first
    placeholder. Body is everything else."""
    placeholders = list(slide.placeholders)
    if not placeholders:
        return None, []

    title = None
    bodies = []
    for ph in placeholders:
        try:
            ph_type = ph.placeholder_format.type
        except Exception:
            ph_type = None
        # Title placeholder types in python-pptx: TITLE (13), CENTER_TITLE (15)
        if title is None and ph_type in (13, 15):
            title = ph
        else:
            bodies.append(ph)

    # Fallback: if no explicit title found, take the first placeholder
    if title is None:
        title = placeholders[0]
        bodies = placeholders[1:]

    return title, bodies


# ─── Universal layout filler (used when agent picks a layout by name) ─────

# Placeholder types we treat as the slide's title.
_TITLE_TYPES = {"TITLE", "CENTER_TITLE"}

# Placeholder types we leave alone when filling — they either inherit content
# from the master (date, footer, slide number, header) or hold an image we
# don't have a source for.
_SKIP_FILL_TYPES = {"DATE", "SLIDE_NUMBER", "FOOTER", "HEADER", "PICTURE", "SLIDE_IMAGE"}


def _ph_type_name(ph):
    """Standard PowerPoint placeholder-type name ('TITLE', 'BODY', ...) or ''."""
    try:
        ph_type = ph.placeholder_format.type
    except Exception:
        return ""
    if ph_type is None:
        return ""
    return getattr(ph_type, "name", "") or ""


def _content_items_from_spec(spec):
    """Return an ordered list of items to drop into a layout's body placeholders.
    Each item is either a string (plain text) or a list of strings (bullets).
    Stops at the first matching shape — agent should not combine forms."""
    # Canonical multi-body form for layouts with 3+ content slots.
    bodies = spec.get("bodies")
    if isinstance(bodies, list):
        return bodies

    # Legacy two-column form (still supported for back-compat).
    if any(k in spec for k in ("leftBullets", "leftTitle", "rightBullets", "rightTitle")):
        items = []
        for prefix in ("left", "right"):
            block = []
            t = spec.get(f"{prefix}Title")
            if t:
                block.append(t)
            block.extend(spec.get(f"{prefix}Bullets") or [])
            if not block:
                items.append("")
            elif len(block) == 1:
                items.append(block[0])
            else:
                items.append(block)
        return items

    if spec.get("bullets"):
        return [spec["bullets"]]
    if spec.get("subtitle"):
        return [spec["subtitle"]]
    return []


def _fill_any_layout(prs, layout, spec):
    """Add a slide using the chosen layout and fill its placeholders with
    whatever the spec provides — independent of any 'type' vocabulary.

    Title goes to the first TITLE/CENTER_TITLE placeholder. Content items go to
    the body/object/subtitle placeholders in the order they appear in the
    layout. Picture/date/footer/slide-number/header placeholders are left
    alone (they inherit from the master).

    If the spec carries more content items than the layout has body slots,
    the overflow is MERGED into the last slot rather than dropped — preserves
    the agent's data when no layout has enough capacity."""
    slide = prs.slides.add_slide(layout)

    title_ph = None
    body_phs = []
    for ph in slide.placeholders:
        type_name = _ph_type_name(ph)
        if title_ph is None and type_name in _TITLE_TYPES:
            title_ph = ph
        elif type_name in _SKIP_FILL_TYPES:
            continue
        else:
            body_phs.append(ph)

    # If no explicit title placeholder AND we have multiple body slots, treat
    # the first body as the title (covers layouts whose title slot is a
    # generic placeholder). For SINGLE-body layouts (e.g. Talentia's "Slide
    # Fin"), don't pop — we'll fold the title into the body content below
    # so the bullets don't get dropped.
    if title_ph is None and len(body_phs) > 1:
        title_ph = body_phs.pop(0)

    title_text = spec.get("title", "")
    if title_ph is not None:
        _set_text(title_ph, title_text)

    items = _content_items_from_spec(spec)

    # Single-body layout with no title slot — prepend the title text into the
    # first body item so it stays visible alongside the bullets/subtitle.
    if title_ph is None and title_text and body_phs:
        if items:
            first = items[0]
            if isinstance(first, list):
                items = [[title_text] + [str(x) for x in first]] + items[1:]
            else:
                items = [[title_text, str(first)]] + items[1:]
        else:
            items = [title_text]

    items = _autodistribute_single_block(items, len(body_phs))
    items = _merge_overflow(items, len(body_phs))
    for ph, item in zip(body_phs, items):
        if isinstance(item, list):
            _set_bullets(ph, item)
        else:
            _set_text(ph, item)


def _merge_overflow(items, capacity):
    """If items exceeds the number of available body slots, fold the overflow
    into the last slot. Lists merge by extending; scalar strings join into a
    flat list. When capacity is 0 or items fits, return unchanged."""
    if capacity <= 0 or len(items) <= capacity:
        return items
    keep = items[:capacity - 1]
    tail = items[capacity - 1:]
    merged = []
    for item in tail:
        if isinstance(item, list):
            merged.extend(str(x) for x in item)
        elif item:
            merged.append(str(item))
    if not merged:
        return keep + [""]
    if len(merged) == 1:
        return keep + [merged[0]]
    return keep + [merged]


def _autodistribute_single_block(items, capacity):
    """Safety net for grid layouts: if the spec arrived as a single block of
    bullets but the layout has multiple body slots, spread the bullets across
    them — one per slot — instead of dumping everything into the first slot
    and letting the empty-placeholder cleanup hollow out the design.

    Triggers only when items is exactly one list of multiple bullets AND the
    layout has more than one body slot. Strings become slot-text (no bullet
    glyphs); a tail list becomes the last slot's bullets when content
    exceeds capacity (the regular overflow merge handles that case
    afterwards)."""
    if capacity <= 1 or len(items) != 1:
        return items
    only = items[0]
    if not isinstance(only, list) or len(only) <= 1:
        return items
    if len(only) <= capacity:
        # One bullet per slot, as plain strings
        return [str(b) for b in only]
    # More bullets than slots — give first capacity-1 slots one each, last
    # slot gets the remaining as a bullet list (so multiple lines render as
    # bullets there).
    head = [str(b) for b in only[: capacity - 1]]
    tail = [str(b) for b in only[capacity - 1 :]]
    return head + [tail]


# ─── Slide builders ───────────────────────────────────────────────────────

def _add_title_slide(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    title, bodies = _title_and_body_placeholders(slide)
    _set_text(title, spec.get("title", ""))
    if bodies:
        _set_text(bodies[0], spec.get("subtitle", ""))


def _add_bullets_slide(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    title, bodies = _title_and_body_placeholders(slide)
    _set_text(title, spec.get("title", ""))
    if bodies:
        _set_bullets(bodies[0], spec.get("bullets") or [])


def _add_two_column_slide(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    title, bodies = _title_and_body_placeholders(slide)
    _set_text(title, spec.get("title", ""))
    # Build left content: optional header + bullets
    left_items = []
    if spec.get("leftTitle"):
        left_items.append(spec["leftTitle"])
    left_items.extend(spec.get("leftBullets") or [])
    right_items = []
    if spec.get("rightTitle"):
        right_items.append(spec["rightTitle"])
    right_items.extend(spec.get("rightBullets") or [])
    if len(bodies) >= 1:
        _set_bullets(bodies[0], left_items)
    if len(bodies) >= 2:
        _set_bullets(bodies[1], right_items)
    # If only one body placeholder is available, fall back to combined
    if len(bodies) == 1 and right_items:
        combined = left_items + ["—"] + right_items
        _set_bullets(bodies[0], combined)


def _add_closing_slide(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    title, bodies = _title_and_body_placeholders(slide)
    _set_text(title, spec.get("title", ""))
    if bodies:
        items = spec.get("bullets") or []
        if items:
            _set_bullets(bodies[0], items)
        else:
            _set_text(bodies[0], spec.get("subtitle", ""))


SLIDE_BUILDERS = {
    "title": _add_title_slide,
    "bullets": _add_bullets_slide,
    "two_column": _add_two_column_slide,
    "closing": _add_closing_slide,
}


# ─── Slide cleanup (pre-load via zipfile manipulation) ───────────────────

def _strip_slides_via_zip(reference_bytes):
    """Return cleaned .pptx bytes with all slide parts removed BUT masters/layouts
    intact. python-pptx's high-level API can't fully delete slides — orphaned XML
    parts stay in the package and collide with new ones, producing duplicates and
    re-rendering of original slides. We solve it by stripping at the zip level
    before python-pptx ever sees the file."""
    output = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(reference_bytes), "r") as zin:
        with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                # Drop every slide part (slideN.xml, slideN.xml.rels, etc.)
                if item.startswith("ppt/slides/"):
                    continue
                content = zin.read(item)

                if item == "ppt/presentation.xml":
                    # Empty out <p:sldIdLst>...</p:sldIdLst>
                    content = re.sub(
                        rb"<p:sldIdLst>.*?</p:sldIdLst>",
                        b"<p:sldIdLst/>",
                        content,
                        flags=re.DOTALL,
                    )
                elif item == "ppt/_rels/presentation.xml.rels":
                    # Drop relationships pointing at slides
                    content = re.sub(
                        rb'<Relationship\s[^>]*?Type="[^"]*?/slide"[^>]*?/>',
                        b"",
                        content,
                    )
                elif item == "[Content_Types].xml":
                    # Drop content-type overrides for slide parts
                    content = re.sub(
                        rb'<Override\s[^>]*?PartName="/ppt/slides/[^"]*?"[^>]*?/>',
                        b"",
                        content,
                    )

                zout.writestr(item, content)
    return output.getvalue()


# ─── Main render ──────────────────────────────────────────────────────────

def render_brand_deck(reference_bytes, slide_spec):
    """Open the reference deck, strip its slides cleanly, add new slides using its
    layouts. Returns the resulting .pptx as bytes."""
    # Pre-strip all slide parts at the zip level before python-pptx loads the file.
    cleaned_bytes = _strip_slides_via_zip(reference_bytes)

    prs = Presentation(io.BytesIO(cleaned_bytes))
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise ValueError("Reference deck has no slide layouts to use.")

    for slide_data in (slide_spec or []):
        # Defensive: some LLMs wrap each slide as a JSON STRING rather than
        # passing it as a JSON object. Parse on the fly so both shapes work.
        if isinstance(slide_data, str):
            try:
                slide_data = json.loads(slide_data)
            except (ValueError, TypeError):
                continue
        if not isinstance(slide_data, dict):
            continue

        # Primary path: agent picked a layout by exact name from the brand's
        # actual layout list (via inspect_brand_layouts).
        layout = _lookup_layout_by_name(slide_data.get("layout_name"), layouts)
        if layout is not None:
            _fill_any_layout(prs, layout, slide_data)
            _remove_empty_picture_placeholders(prs.slides[-1])
            _remove_unfilled_text_placeholders(prs.slides[-1])
            continue

        # Fallback path: agent didn't supply a layout_name, or the name didn't
        # match anything. Use the legacy heuristic + type-specific builder.
        slide_type = (slide_data.get("type") or "bullets").lower()
        layout = _pick_layout(slide_type, layouts)
        if layout is None:
            continue
        builder = SLIDE_BUILDERS.get(slide_type, _add_bullets_slide)
        builder(prs, layout, slide_data)
        _remove_empty_picture_placeholders(prs.slides[-1])
        _remove_unfilled_text_placeholders(prs.slides[-1])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─── Orchestrate tool wrapper ─────────────────────────────────────────────

@tool(
    name="brand_render",
    description=(
        "Renders a presentation by adding new slides INTO the user's reference .pptx "
        "using its slide masters/layouts. The output automatically inherits the "
        "reference's logos, colors, fonts, footer, and accent shapes — no separate "
        "style extraction needed. "
        "Inputs: pptx_file (the reference .pptx the user uploaded) and slides (a list "
        "of slide dicts). "
        "PREFERRED: each slide includes 'layout_name' chosen from the output of "
        "inspect_brand_layouts (call that tool first). The renderer fills the chosen "
        "layout's placeholders in order: 'title' goes to the title placeholder; "
        "content goes to body/object/subtitle placeholders sourced from 'bodies' (a "
        "list whose entries are strings or string-lists), or from 'bullets', or from "
        "'subtitle'. An optional 'layout_reason' field is accepted and ignored at "
        "render time (kept for the conversation log). "
        "FALLBACK: if 'layout_name' is missing or unknown, the renderer falls back to "
        "a small heuristic that uses the legacy 'type' field — one of 'title', "
        "'bullets', 'two_column', 'closing' — with associated fields ('subtitle', "
        "'bullets', 'leftTitle/leftBullets/rightTitle/rightBullets'). "
        "Returns the rendered .pptx as a downloadable file."
    ),
)
def brand_render(pptx_file: WXOFile, slides: list) -> bytes:
    reference_bytes = WXOFile.get_content(pptx_file)
    return render_brand_deck(reference_bytes, slides)


# ─── CLI for local testing ────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("reference_pptx", help="Path to the brand's reference .pptx")
    parser.add_argument("spec_json", help="JSON file with a 'slides' array")
    parser.add_argument("output_pptx", help="Where to write the rendered .pptx")
    args = parser.parse_args()

    ref_bytes = Path(args.reference_pptx).read_bytes()
    spec = json.loads(Path(args.spec_json).read_text(encoding="utf-8"))
    slides = spec["slides"] if isinstance(spec, dict) and "slides" in spec else spec
    out_bytes = render_brand_deck(ref_bytes, slides)
    out_path = Path(args.output_pptx)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_bytes(out_bytes)
    print(f"Wrote {out_path} ({len(out_bytes):,} bytes)")
