"""Diagnostic: list every slide layout in a .pptx with placeholder details.
Helps us tune brand_render's layout-picking heuristic for a specific deck.

Run:
    python prototype/inspect_layouts.py "path/to/deck.pptx"
"""
import sys
from pptx import Presentation


PLACEHOLDER_TYPE_NAMES = {
    13: "TITLE",
    14: "BODY",
    15: "CENTER_TITLE",
    16: "SUBTITLE",
    17: "DATE",
    18: "SLIDE_NUMBER",
    19: "FOOTER",
    20: "HEADER",
    21: "OBJECT",
    22: "CHART",
    23: "TABLE",
    24: "CLIP_ART",
    25: "DIAGRAM",
    26: "MEDIA_CLIP",
    27: "SLIDE_IMAGE",
    28: "PICTURE",
}


def main(path):
    prs = Presentation(path)
    print(f"=== Slide layouts in {path} ===\n")
    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name or "(unnamed)"
        ph_count = len(list(layout.placeholders))
        print(f"Layout {i:>2}: '{name}'  ({ph_count} placeholders)")
        for ph in layout.placeholders:
            try:
                idx = ph.placeholder_format.idx
                t = ph.placeholder_format.type
                t_name = PLACEHOLDER_TYPE_NAMES.get(int(t), f"type={int(t)}")
            except Exception:
                idx = "?"
                t_name = "(unknown)"
            ph_name = ph.name or "(unnamed)"
            print(f"           idx={idx} type={t_name:<14} name='{ph_name}'")
        print()


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python prototype/inspect_layouts.py <path/to/deck.pptx>")
        sys.exit(1)
    main(sys.argv[1])
