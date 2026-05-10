"""Talentia-branded PowerPoint renderer (Stage 1: layout + typography, no logo yet).

Reads the same JSON shape as siro_simple_render.js so it's drop-in compatible
with the existing agent's deck construction:
  { "title", "subtitle", "settings": {...}, "slides": [ ... ] }

Run locally:
    python prototype/talentia_render.py prototype/sample_deck.json prototype/output.pptx
"""
import argparse
import base64
import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Slide dimensions — 16:9 widescreen, matches LAYOUT_WIDE in siro_simple_render
SLIDE_W = 13.333
SLIDE_H = 7.5

# Talentia defaults; everything is overridable via the input JSON's settings
DEFAULTS = {
    "fontFace": "Arial",
    "titleColor": "24135F",       # Talentia blue
    "bodyColor": "1D1D1D",
    "backgroundColor": "FFFFFF",
    "accentColor": "D0006F",      # Talentia magenta
    "footerText": "Generated presentation",
}

LIGHT_PANEL = "F5F5F5"
WHITE = "FFFFFF"


# ─── Settings helpers ─────────────────────────────────────────────────────

def _clean_hex(value, fallback):
    if isinstance(value, str):
        v = value.lstrip("#").upper()
        if len(v) == 6 and all(c in "0123456789ABCDEF" for c in v):
            return v
    return fallback


def _hex_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _normalize_settings(raw):
    raw = raw or {}
    out = {}
    for key, default in DEFAULTS.items():
        if key.endswith("Color"):
            out[key] = _clean_hex(raw.get(key), default)
        else:
            v = raw.get(key)
            out[key] = v if isinstance(v, str) and v else default
    return out


# ─── Drawing helpers ──────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill_hex):
    """Filled rectangle with no visible border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _text(slide, text, x, y, w, h, *,
          font="Arial", size=14, bold=False, italic=False,
          color="1D1D1D", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_rgb(color)
    return tb


def _bullets(slide, items, x, y, w, h, *, font, size, body_color, bullet_color):
    items = items or []
    if not items:
        return
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        marker = p.add_run()
        marker.text = "▪  "
        marker.font.name = font
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = _hex_rgb(bullet_color)
        body = p.add_run()
        body.text = str(item)
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = _hex_rgb(body_color)


def _add_logo(slide, x, y, width):
    """Drop the embedded Talentia primary logo at (x, y); height auto-scales from the PNG."""
    img_stream = io.BytesIO(base64.b64decode(LOGO_PRIMARY_B64))
    slide.shapes.add_picture(img_stream, Inches(x), Inches(y), width=Inches(width))


def _footer_band(slide, page_num, settings):
    _rect(slide, 0, 7.1, SLIDE_W, 0.4, settings["titleColor"])
    _text(slide, settings["footerText"], 0.5, 7.18, 9, 0.25,
          font=settings["fontFace"], size=9, color=WHITE,
          anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, str(page_num), 12.5, 7.18, 0.5, 0.25,
          font=settings["fontFace"], size=10, bold=True,
          color=settings["accentColor"], align=PP_ALIGN.RIGHT,
          anchor=MSO_ANCHOR.MIDDLE)


# ─── Slide builders ───────────────────────────────────────────────────────

def _title_slide(slide, spec, s):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, s["backgroundColor"])
    # Top accent bar (thin band of accent color across the full width)
    _rect(slide, 0, 0, SLIDE_W, 0.22, s["accentColor"])
    # No logo — extraction proved unreliable across decks (most brands draw their
    # wordmark as text on the master, not as a PNG asset).
    # Big title
    _text(slide, spec.get("title", "Untitled"), 1.0, 2.4, 11.3, 1.7,
          font=s["fontFace"], size=44, bold=True, color=s["titleColor"])
    # Short accent rule below title
    _rect(slide, 1.0, 4.15, 1.5, 0.07, s["accentColor"])
    # Subtitle
    if spec.get("subtitle"):
        _text(slide, spec["subtitle"], 1.0, 4.45, 11.0, 0.8,
              font=s["fontFace"], size=22, color=s["titleColor"])
    # Bottom-left accent corner mark
    _rect(slide, 0, 6.9, 0.6, 0.6, s["accentColor"])


def _bullets_slide(slide, spec, s, page_num):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, s["backgroundColor"])
    # Left vertical accent bar (stops before footer)
    _rect(slide, 0, 0, 0.22, SLIDE_H - 0.4, s["titleColor"])
    # Title
    _text(slide, spec.get("title", "Slide"), 0.65, 0.55, 11.5, 0.7,
          font=s["fontFace"], size=28, bold=True, color=s["titleColor"])
    # Thin accent underline below title
    _rect(slide, 0.65, 1.32, 4.0, 0.05, s["accentColor"])
    # Bullets
    _bullets(slide, spec.get("bullets"), 0.75, 1.7, 11.5, 5.0,
             font=s["fontFace"], size=18,
             body_color=s["bodyColor"], bullet_color=s["accentColor"])
    _footer_band(slide, page_num, s)


def _two_column_slide(slide, spec, s, page_num):
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, s["backgroundColor"])
    _rect(slide, 0, 0, 0.22, SLIDE_H - 0.4, s["titleColor"])
    _text(slide, spec.get("title", "Comparison"), 0.65, 0.55, 11.5, 0.7,
          font=s["fontFace"], size=28, bold=True, color=s["titleColor"])
    _rect(slide, 0.65, 1.32, 4.0, 0.05, s["accentColor"])

    card_y = 1.75
    card_w = 5.6
    card_h = 4.5
    left_x = 0.65
    right_x = left_x + card_w + 0.5

    # Left card
    _rect(slide, left_x, card_y, card_w, card_h, LIGHT_PANEL)
    _rect(slide, left_x, card_y, card_w, 0.55, s["titleColor"])
    _text(slide, spec.get("leftTitle", "Option A"),
          left_x + 0.25, card_y + 0.10, card_w - 0.5, 0.35,
          font=s["fontFace"], size=18, bold=True,
          color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _bullets(slide, spec.get("leftBullets"),
             left_x + 0.3, card_y + 0.8, card_w - 0.6, card_h - 1.0,
             font=s["fontFace"], size=14,
             body_color=s["bodyColor"], bullet_color=s["accentColor"])
    # Right card
    _rect(slide, right_x, card_y, card_w, card_h, LIGHT_PANEL)
    _rect(slide, right_x, card_y, card_w, 0.55, s["titleColor"])
    _text(slide, spec.get("rightTitle", "Option B"),
          right_x + 0.25, card_y + 0.10, card_w - 0.5, 0.35,
          font=s["fontFace"], size=18, bold=True,
          color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _bullets(slide, spec.get("rightBullets"),
             right_x + 0.3, card_y + 0.8, card_w - 0.6, card_h - 1.0,
             font=s["fontFace"], size=14,
             body_color=s["bodyColor"], bullet_color=s["accentColor"])

    _footer_band(slide, page_num, s)


def _closing_slide(slide, spec, s):
    # Solid primary-color background
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, s["titleColor"])
    # Big centered white title
    _text(slide, spec.get("title", "Thank you"), 0, 2.5, SLIDE_W, 1.4,
          font=s["fontFace"], size=44, bold=True, color=WHITE,
          align=PP_ALIGN.CENTER)
    # Accent rule centered
    _rect(slide, SLIDE_W / 2 - 0.75, 3.9, 1.5, 0.07, s["accentColor"])
    # Centered white bullets (no markers — cleaner look at this size)
    bullets = spec.get("bullets") or []
    if bullets:
        tb = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.4), Inches(SLIDE_W - 4.0), Inches(2.2)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = str(item)
            run.font.name = s["fontFace"]
            run.font.size = Pt(20)
            run.font.color.rgb = _hex_rgb(WHITE)
    # Bottom accent strip
    _rect(slide, 0, 7.2, SLIDE_W, 0.3, s["accentColor"])


# ─── Driver ───────────────────────────────────────────────────────────────

def render(input_path, output_path):
    deck = json.loads(Path(input_path).read_text(encoding="utf-8"))
    settings = _normalize_settings(deck.get("settings"))
    slides_spec = deck.get("slides") or []

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank_layout = prs.slide_layouts[6]

    for i, spec in enumerate(slides_spec):
        slide = prs.slides.add_slide(blank_layout)
        slide_type = (spec.get("type") or "bullets").lower()
        page_num = i + 1
        if slide_type == "title":
            _title_slide(slide, spec, settings)
        elif slide_type == "two_column":
            _two_column_slide(slide, spec, settings, page_num)
        elif slide_type == "closing":
            _closing_slide(slide, spec, settings)
        else:
            _bullets_slide(slide, spec, settings, page_num)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


# --- Embedded assets ---
# Base64-encoded Talentia primary logo (PNG). Keeps the tool self-contained.
LOGO_PRIMARY_B64 = "iVBORw0KGgoAAAANSUhEUgAAAhUAAADSCAYAAAD9ljJbAAAgp0lEQVR4Xu2dP8s9SXbf90XoDdgvQODc4FyBYidSaCVWYhQZRQ4MyhQKwWYGBQLDBlrsTRbMglkwmyxCm2ygUbDBJutgk01GfGeo2TNnuv6eU3X79v184Mv85nm6q6qrTp36dnXf+3zvSwAAAIAEvud/AAAAALACpgIAAABSwFQAAABACpgKAAAASAFTAQAAAClgKgAAACAFTAUAAACkgKkAAACAFDAVAAAAkAKmAgAAAFLAVAAAAEAKmAoAAABIAVMBAAAAKWAqAAAAIAVMBQAAAKSAqQAAAIAUMBUAAACQAqYCAAAAUsBUAAAAQAqYCgAAAEgBUwEAAAApYCoAAAAgBUwFAAAApICpAAAAgBQwFQAAAJACpgIAAABSwFQAAABACpgKAAAASAFTAQAAXf7xe3+eKngmmAoAAOjiTUFU8EwwFQAA0MWbgqjgmWAqAACgizcFUcEzwVQAAEAXbwqigmeCqQAAgC7eFEQFzwRTAQAAXbwpiAqeCaYCAAC6eFMQFTwTTAUAAHTxpiAqeCaYCgAA6OJNQVTwTDAVAADQxZuCqOCZYCoAAKCLNwVRwTPBVAAAQBdvCqKCZ4KpAACALt4URAXPBFMBAABdvCmICp4JpgIAALp4UxAVPBNMBQAAdPGmICp4JpgKAADo4k1BVPBMMBUAANDFm4Ko4JlgKgAAoIs3BVHBMzluKv7vT/7py3/zB3/6sfrrv/qfvkvgAt9vUQFE0dz1cfVJUu4G6IGpOCxMxRi+36ICiIKpwFRAH0zFYWEqxvD9FhVAFEwFpgL6YCoOC1Mxhu+3qACiYCowFdAHU3FYmIoxfL9FBRAFU4GpgD6YisPCVIzh+y0qgCiYCkwF9MFUHBamYgzfb1EBRMFUYCqgD6bisDAVY/h+iwogCqYCUwF9MBWHhakYw/dbVABRMBWYCuiDqTgsTMUYvt+iAoiCqcBUQB9MxWFhKsbw/RYVQBRMBaYC+mAqDgtTMYbvt6gAomAqMBXQB1NxWJiKMXy/RQUQBVOBqYA+mIrDwlSM4fstKoAomApMBfTBVBwWpmIM329RAUTBVGAqoA+m4rAwFWP4fosKIAqmAlMBfTAVh4WpGMP3W1QAUTAVmArog6k4LEzFGL7fogKIgqnAVECf46biX7749VeTM1v//t/9l+9Mgqh8HRliYo7hxyIqgCiau34+Z8jHalTKhb6ODCl3A/Q4bip28R//+L9/Z3JFBa/Dj0VUAHfFx2pUyoUArwJT0RC8Dj8WUQHcFR+rUWEq4JVgKhqC1+HHIir4LPSo4l22632sRoWpgFeCqWjojihZSn//d//nm2edunbJPv/U7//x5//sT38b/FhEdVf+/29+ezmm/+lP/vqrMf1v//V/fPOzctw7j2smMg3qs7/4z3/bnf9/9B/+8qtjdOz//of/91W/3wXf1qgwFWvYufj9v/lfX825Elv6r82viqF3Ma2nwVQ0dAcU6EqcWmR8+0ZVDEdtMbKTKUMZk81fQ1R3QglJZiHycrHOVaJTbNxhgcyOoVasKpZlEnyfzErzQotHRrxG8O2Kapep8PVEdQc09oqB1XjSPFRubsWRYtnHd0R3mO8tMBUNvZKokahJk8dPAAWqPy4iJf0ovsyoXomSgPp8x3gWaVzV769KONkx5BfG8qmxP/y3f/adY6NSma/sO9+eqHzfZeHriepVKJZk6leNREua48rdluy1SXPtzmAqGnoFuouN3MHOqCTS7AUBU/E15a56x0JY06sWyOwYsgujrsf/fode1Xe+HVFhKq45OR+Vw4u5yF6bMBWHyB446SQKlB3X0JMmWPYdNKbi652mU+bwSqcXyB2mQneUO+4me1Kdta3sHfj6o8JUfBfNxxNmwks5IDsPYCoOsWNBPoGSvrbifN3vrE82FSd3mkakttTeT8gk21RoAXjFIlCkuk/0m/B1R4Wp+D0aw1cY053CVBziHU2FDMXTAl76VFNxZ3Pon/Nmk20q7iAZi939Jny9UWEqvkYx+UpjukuYikO8m6mQg35iwEufZirexRzq0yK7Hoc80VQUafdpJ76+qDAVXz/u8PU9RZiKQ7yTqVCSeqqhkD7JVLyLoShSW3cYiyebCs3Vne9Y+Pqi+nRTcecdwwxhKg7xLqbiycm36FNMxbsZiqIdxuLpca0+24WvK6pPNhVPNxQSpuIQ72AqdLfz5B2Kok8wFe9qKIqyF56nmwpJC9YOfD1RZY9twdcTVTZPfuRhhak4xN1NxbsvQjP6BFOR/THcVyhjnAqfYCqkHY9BfB1RfaKpePI7al6YikPc3VToJTlf/lOVsVj5MqPKRN+O6ct/V2W9hPgppkLzOBtfR1SfZip0w3anj3HvFqbiEHc2FU9ahEb0ZFPxtEdYSsYZfIqpkLJ3K3z5UX2aqdiR++8sTMUhdgRWBnLRT1qERvRkU/HEHaeM72L4JFORvVvhy4/qk0zFJ8VdEabiEHc1Faf+bsGd9FRToTtUX262ZEAVy0Un3sPRbkX00yCflNw1RtH+svjyo/okU7Ej799dmIpD7AiuKCd3KbQwqA90F6VFXSoLkz92t55qKnbtUihG1Getr4XW7/T+w67xjI7ZJ5kKKetdFOHLjupTTMXJT3vI3Ktf9QkgzRX995Tp98JUHGJHso2ye5dCi5GCu7UYFZQEtSieMDnRBUr4MqOKoj72ZWZIfTV717vj7xlE775faSpe8ZJe5sdLfdlRfYqp2D3uKl/vw/XeodHvZXB2rEFXwlQcYseARti9SyGD0Av2K3TOrjvuoieaiuyPkCo2Rsxgi+w2Rd6tOGkqNNfV1qv4Vzu04O+ce1Lml2H5sqP6BFOxc5ei7ByuoPjLNvxemIpD3M1U7Az6SPIvKDB3Jd7VCWnxZUYVxZcXUYahENnffRJ5AfGEqZCJujISV6hvdpvnyM6OxZcb1SeYimxDXaT5lDE3d37iD1NxiLuZih1Bn7UYFZSgMxeloqeZCj068uVFlPk8PvMjripnlZ2mQu1aNdI7zX3WOPpyo/oEU+HLylD219fvij1MxSHuZip8WRnKSmKW7Ltd6WmmIvPvCchsZpP57s6qad1lKhSbo7sTNTL7xyojzoUvN6qnm4psky9F3ymqseObPjEVh7iTqdgR9JkvhnmyF4SMZOvLjCpC5gthOxJC5kddV+MsO4akzDvHbOMsZcS58OVG9XRTseOx1o55Wchem3a2NQNMRUOrZN7ZSrtcdCF7QchItr7MqFbJXLCzvr3yiqzHbasLUnYMZXx3hmXHVvSqAfP4cqNaHcMevp6oVsk0+VLkXaIRstcmTMUhsgdOWiU76DMW6RbZC0JGe32ZUa2S+cKVytpF1hb/qvHJjqEdC2P2vMxqoy83qqx2eXw9Ua2w46Pd0cdrPbLXJkzFIbIHTlrFlxPR7l0Kkb0gPMlUZC3WUvkCnR3K3N5fITuGdE3ZZG+bZ7XRlxtVVrs8vp6oVsiOsx3vOHk0Hr7eiDAVh8geOGkVX05EJ4I+e6I+yVRkP8p6B62QHUM7FsbMXScpq42+3Kiy2uXx9US1QvZjrNVPFc2QvTZhKg6RPXDSCtnJdeeWeSG7zU8yFTvi6u5aSVrZMbRjYbxrG325UWW1y+PriWqFzJ1DafejD5GdQ1bm50kwFQ2tkJ24TgRQdpsxFe+tlZjLjqEdC+Nd2+jLjSqrXR5fT1QrZJqKyPeyzJCdQ1bm50kwFQ2tkP1x0t3vU4jsZPskU5H9ct87aCVpZcfQjoVRc8nXE1FWG325UWW1y+PriWqFzDy/q588mW2WVubnSTAVDa2Q6aSlE2QvCE8yFb6cT9BK0sqOoV0J39cTUVYbfblRZbXL4+uJaoXMPL+rnzyZbZZW5udJMBUNrYCpwFS8u1aSVnYM7Ur4vp6Istroy40qq10eX09UK2Tm+V395Mlss7QyP0+CqWhohezHHydeJMpeEJ5kKnj8MUZ2DO1K+L6eiLLa6MuNKqtdHl9PVCtk5vld/eTJbLO0Mj9PgqloaIXs5HoigLLb/CRTsSOu7q6VmMuOoV0J39cTUVYbfblRZbXL4+uJaoXMnWBe1NwDpqKhFbKT6/f5SGmKVtkRV3fXysvB2TH0DgtjVht9uVFltcvj64lqhUxTIZ3YCc7OIZiKQ2QPnLSKLycivvwqR6vw5VdjZMfQOyyMWW305UaV1S6PryeqFfjyK0zFMbIHTlrFlxPVyp3jDNkLwpNMRfad0TtohewYeoeFMauNvtyostrl8fVEtUJ2nJ24actemzAVh8geOGmV7Jf7MhbpFtkTNaO9vsyoVtHjJ1/WqnYl+zuQHUO7+srXE1FWG325UWW1y+PriWoF/qAYpuIY2QMnrZK9Zb77j4plLwhPMhWZf/pc2p3AXkV2DL3DwpjVRl9uVFnt8vh6olol+6aNP32eC6aioVWyP1Yq7Qz87AXhSaZCZCaxneP4SrJj6B0Wxqw2+nKjymqXx9cT1SrZN22ScvYustcmTMUhsgdOiqDdBV9eVDteKtIOSOafzZaeZiqyk5i2cJ8GpmIdX25UWe3y+HqiWmXHTZvy9Y5dRM2L7LUAU3GIu5kKvQDky4tKwZnpqLW4ZRsK6WmmInvBVJ9nPs7SOKqNGVo1PDrXX2dE77AwZrXRlxtVVrs8vp6oImQv1JLmZaaxyP6kSpHm2p3BVDQUYVdASRk7FgrMHRNTepqpENl9pXjNMBbZcbb6eAZTsY4vN6qsdnl8PVFF2HHTJmmerxpry/cTX/D2wlQc4m6mQgtG9kJkpetdCX458V0TsuiJpkKLrS8zquidUfZCLq3uhGW35R0Wxqw2+nKjymqXx9cTVYRsM+2lR54rpl/zIPMdrCthKg5xN1MhTnzHgQyCJlhrAmjh0jG7zUTRE01F9qdAimQ81V+t8fPo2Oz3PKTI1xZjKtbx5UaV1S6PryeqKDse3VppPmie9Yy2bu60M7G7PUWYikPc0VTs3q3wUlCrH6x2u+YrPdFUiB27FVYqX8npagdKpkbJTcfsiqnVRx8CU7GOLzeqrHZ5fD1RRcmOuZ58bpV2zcWWMBWHuKOpECd2K+6mp5qKXbsVd1EkWWUn+HdYGLPa6MuNKqtdHl9PVBnsyPt3V2SengBT0VAG2q14xW7BK/VUUyGeahKjCxGmYh1fblRZ7fL4eqLKIDvu3kGYikPc1VSIHZ+rvrOebCpOP9I6patHLjNkJ/d3WBiz2ujLjSqrXR5fT1RZnHpX7C7CVBzizqZC7H4efyc92VSIp5nEyLsUBUzFOr7cqLLa5fH1RJXFji/wu7MwFYe4u6kQnxL4TzcVYsenL14h7brMfPKkBqZiHV9uVFnt8vh6ospEO21P3EG8EqbiEO9gKp66de71CaZCvPu2a9YX/QhMxTq+3Kiy2uXx9USVzdN2EGvCVBziHUyF2PGne++mTzEV777tmvHNrAVMxTq+3Kiy2uXx9US1g6e+SG2FqTjEu5gKoaB48o7Fp5gK8a7GItNQCEzFOr7cqLLa5fH1RLWLpxsLTMUh3slUiF1/zOsO+iRTUXiXF3FlZrMNhcBUrOPLjSqrXR5fT1Q72f013q8UpuIQ72YqxLve5fb0iaZCfH/jHxHKUOY7FB5MxTq+3Kiy2uXx9US1m6e+vImpOMQ7morCUz5JUPSppkLcdQdK8yPjUx41MBXr+HKjymqXx9cT1Qn0Lbg71oZXClNxiB2BcxItRjuuoSc5+exPMXyyqSho+/UO36SqNpxIQpiKdXy5UWW1y+Prieokr5qPusHIrvfEfI6AqWjoFShgsoOwJi3+unvNXhAwFb/nVclMdepxzCmyY+gdFsasNvpyo8pql8fXE9UrUG468UhE86+8u5S9NmEqDpE9cNIr2fWnyuWctdhoW7CQvSBgKr5LGc/dCU0vjPb+VPMOsmPoHRbGrDb6cqPKapfH1xPVq9CNlHLgjseUmuP+RejstQlTcQgNpBazTN0BTQBd2+qfvNY5CmpdT+0lPRkMf+0RZQS9v46o7oT6R/2UkdRUhmJDMbLznYke2THkE3MWvp6Istroy40qq12fgOI2YjC0IyEjoTJq8y97bbI3hHfkMabiU5Ax0KJkA1WmQYFtA0+/r5kIuA/l8VMxGnZMvezvM4wbAHybMhdlEjTP9BK95p6Mu51/2g2smYhPB1MBAAAAKWAqAAAAIAVMBQAAAKSAqQAAAIAUMBUAAACQAqYCAAAAUsBUAAAAQAqYCgAAAEgBUwEAAAApYCoAAAAgBUwFwAF+97vfffmb3/zG/xgA4FE8xlT8/Oc///InP/nJkq74wQ9+8I1alDK++OIL/6suOqfVhhF++ctfflPGr3+99odmtNiVMlYWvnKuxmAW1afzfvzjH3+rz3/0ox99VeYvfvGLL3/727nv2M+OhVVkJH72s599+cMf/vBb16Zr/dWvfuUPT6Vcz+yY2FhoxbSN3ZpUt8ZvNS5nKO1R3Pi+1hisxHVB46h59tOf/vRbZZc4Vfkr12j7MNI+MTJmNex88YyM84iu2jWaY69QTlBs+bxRxlzXtNKnpS9m541YyaO6DsWPzrF5Qv/Wz3SNir934jGmQgPgg2tUV/R+X7DHjQZSQQEzUkcLm0TVBysoIZYyNCFng3ilfk2mqyRd08zkyo6FFdSn3kx4KZnsotQxMybCxoL6vIaN3REpTq8WlihqrzcSNakvZg2qrrM3jrb8GXNh+3DmvCtKOa0xq2Hni2d2nGu6apf9/SjFqPvya9KxM2Nu+0JGcgY7d0bGU8bFt/dKir8Vk/MqHmkqilsc1RV2UFv4wZ8J4KipUJL2AThTf8FOhtJ/M8yep7t0n6itM5dkOPxioWNGjFt2LMyixGevT9eia1Y/a8zs72YT1yizY1JYMRUyor4fJTt29thRc9hDsWD7UvGi5Ks+1nWob/X//g5wpH4dU7sL1kKl69c1+jiWRs3TO5iK3k6F7Vf/O6urPrF9NoLG2+cESeWXXTHbJtu20f7154+eJ2ZMhTdrqlc/03nKFVfXsvMmJJNHmooM7GC2sMdJM0kzairKBLOJbSXwvKmQZpxxOUdj0MPXpXNaE1DJyCaSEeOWHQuz2HG9SvLWdOi/O5gZE8uKqWiNn65Vx9oYnYmtFjYuWnHv725l8lp4Q6G2t8zf1a5b6/jCaB+OUMppjVmNyHwp583GmbD91ePKQLYeIao//aI80sf+nJF8Uxg1FboWW37rWP3OXnfrmu8CpqKCDawW9riiVoKzREyFT/72+kcnQcEv9EVXdxdXlON7iUXtWllcfJLv1ZMdC7OM1K9FpxzTSiqrjPaVx8dVjdkF0S8Ko8a7hmKplKXY6KH6rAlp1W/HT2WPzie/c9hbAGb7sEUppzVmNUbitUY5bzbOhO2rFn7slF9b42exZnLEIHhTIY3eKI6aCjvuI+bTmpDRteWVYCoq2KBqUY5R0NukORIsEVPhTYQSWPn/2cRiJ4O/Qxt53FCO7yUWO8FnJ4dPLK3kkB0Ls5R26r811H61Uxo1bzOMjolnl6kQK+fUWIl3u+jX+tyWq/EbWUwsto7W+IvM/pjtC0tkvpTzZuNMlHN79dp+WqlnJu94Qzl6nhg1FXZXq5XHLKUtu3Y2M8FUVChl9corx6h+m5B6gSVWTUXNudrHITPJ0C8kdhKOJFbbBzVsm3vJtoZN2K1djuxYmMUmjVcxMiZX+FiosbIg2hholT2CnWutWLAojtVWqZbMrXEdvS6PHf+aeRErfVgj0q+R+VLOm40zUc5t1asxK8fM5jWLHdfWjZLtC7+z2rtRHDUVto5WWyxqS4nd1T44BaaiQimrV145pkwqmyh6222rpsIu+jYo7aLbmwCWq4Vk5nHDyHH2WluJtkeZ5C3Hnh0Ls9hHGytJPoORMbniKhauWFkQM02Ffy7dmmej2DJn+81iH8203t9Y6cMakX6NzJdIf5VzW/WO3kj0GC3H94WdD1LrkdaoqbDj3oqPdwVTUcEGUotyjJ1UdtFvPY9bMRU2YfmJrHrKojuzG3C1kPjHDa2JWGuPZcWdr5IdC7P4vtt9vVeMjMkVV7FwxcqCOPL4YQb/qE5l1ubaCJlm0N7h1ljpwxqRdkfmSzlvNs5EObdVr931ifRRK29arvrCxkXrkfCoqbDmtbSnZVbeDUxFBTvoLa4CVYlt5Hnciqmw51wFov39aOKuLST+5bpaeVd94BlJslnYWNC1jaiWKFaxC6iu/WqsdlLqVhz6a21pdGGdXRCt4ZUy+tvHp71mzTm1caRtBfu9ATPnXTFiomf7sMXImNWI5M5yXmvu17BjVsO2LWIYxUgOqvXFyI2ixrAc0xvPq++oUPtUf/lYdMbu2yt4pKnQv0c0cvftg8tj67SMPI+bNRUjOxGjjtxiJ4NPSn5xvEqQI/WNHJOFjYVRZbbLjqtVzZTtwNe9Ih8LlpkFUYZqxGSvUF549W23UtzqjrfXTlvO1aIxw0j/jBwzSimnNWY17HXPUs5bmT92jGr08t0MI9dZO8bfKF49tpgxFeLq+3q8dN1ap6LxeJJHmopRtSaCPa5FqywbZJIPtFlTYe8iWwuUddW+zitapkJYV3314marDwrlGE3M3WTHwgz+I2z2/6XWuI3uEozgr29FrTbY2C27Al7qU/sYSNL477gDU5nqv6svTbPS7338Flbeyq8xsuuBqfj9uNSwL59HsaagRqsv/G6b7+tZU1HQsYqXVt5Sva3ccSceaSr07xHt3Kko+Dt9m6xmTcXopztscF85ak/PVAjfv5bazy020e/GtlXXNqKrHZhZ7FjbLVJ/R3KVHOzOVobxsu3w19rSqLGxsTsqjUsrbjNRPboezXG7mEi1eW+vKfq4ysZgzaDsMBUru0C2rbPYsZ3FjkmNkX4cpZTTMii9vvDvQ9g40RiWn0fGU3UoR2gs/U5GRp7azSNNRQZ2IFuUY1qTqvY8bsZUeHNizdGVbPt7k9FOhtpC4rf/bPIa6YOsF640qcoCWCM7FkZpmT7//N8n/9K33niuMjImV4zEgpgxFWpDxjVFsGZJukrOdo61rn2EUk5rAcs0FSW2ZsdblNhbecRgx3gWOx41bO6MGD1rBlo3WiO5w+fiEktZpsKiPGJz50o/nwZTUaGU1StvdLDtglyCesZUeKMwI7+AeUYXEr8wljvu8v+tPshI2JpgpYxWAsyOhVFKnbU7YfWf3bEp42Lj4OrdmxVGxuSK0VjoLYi1xJuBDIrql67qrmEXqKvdotH46mGvvTX3en04QyTmy3mzsSIyzm212X4fSasve4zOsdF+vLpRHDEVio0Su6Oo7JM7vVEwFRVKWb3yyjG9SXX1PG7UVNhgVT3l3J5sIPo7Z8voQiL8F3zZO4BWH9jrX70bz04MmYz2od/xsW1t9d8sq2WOXocdi1oCtcdcvYuzil38W3ednpFrs+PRirEafgGo9Y0Y6cNR7DscM3f0UbNfzp2NM1HOlVqM9mcN+2jxahfRMpM77LHlReBeO+2O2cw4zbTr1WAqKpSyeuWVY0YmlX8eZxeXFquPDkYTxkiytdhkaI1Srw/sefYx0Ah+l6R1bnYsjFLq7N1ReWNR+nHFaNUYHRPPaCyMLoi7tm5t/7ViwWITeq3N9tNT0uwOi72D7V3vaB+OYOf6qIHzNzorbRi91itsP7fwu14j11bwc60V02Imd3gDac+t9eXoYxhPuYbW47S7gKmoUMrqlWcDagQ7QUbqsElutI6C3x2oMbqQWPynGkba5ye4/j2ykHpDUXu8UMiOhVHsxO8tRv6alJxG+mKU0THxjMbC6ILox7xnuEaxBkHl9xYa347W8f7Ljq4elXhUnjVQIyZxtA9H8Y/WWtfo2zsbJ4XI+eVcqYfv2978Eup/n29afSJmc4efx0Wt8bTjNLIbZteMGSPyKjAVFWyAtCjHzEyqqy8+qdF7DtzD1lU7f3QhsfgkPdoH/u5I/1adV5NQScGbF9XZIzsWRvF3VLXtTSWSq0Q0mixHmBkTy2gszCyIPvHW4nAGf5eof9fKVftmjY2PO/WjxtMbBZWtvvAxPTKOM304gt8JVZ+ozXYh1b/9p5FG23uF7Z9ZbFt7XOUbjZH6zV+ffuZz7OjuzUruuLpRbI2nf4Ss6/BxJdRenyta5d4FTEUFO+gtyjGzk8q2t1aHgqr8fvXFMbvTUStjdCHx2GeVM31wlSCK9HPfN7b82cQwqyh+MdL16O6i9jl0/dzfGa8meEspb3RMCqOxMLsg2nKljGv0ZsVec5E1HtKIKS3Ya7RSnbX4VX2j11Yrv6fWmGqBu+oTtcv3hdQyvyOMtKmGbccImvt+fvWur7RtJG+I1XXEj2VvTlwZEY2FjV3/+5FdjTuAqahgB7NFOWZ2Uvk7rStsoLYSfA87Ea8SyOhCcoU9d7YPNLFqicBK5fYmqeVqQo4qA3/neqVyF1nwSaZ21z1KKWd2TEZjYdZUCG+eRhN9C5Ux0t/6vY6brVMGYSSeVsr3C9GoemM62mYdM2qAaoy26QrblhkUbyN5o7V7VSOyjsy++6ZjRsbJ54q78xhToeDRJG0lwh2UOmeDV2hCt9psr2kmWXlsPVfBvvoRvYICfrUPhOrU+dallzv4lfbYfptVFmXr0l+XDF4tQZR+zGhLKWN2TEZjoYyZdLV1W6P0Sa/8WdTfZdyV3NXX+m/pg8j8EbpGtV3jZ8dT5dfGs4ftwxmNjqnmfYkp396omSjMtikTXYMfkzLmq9cXWUeKwZVm5oTNzzZPlLF6Nx5jKgAAAOC1YCoAAAAgBUwFAAAApICpAAAAgBQwFQAAAJACpgIAAABSwFQAAABACpgKAAAASAFTAQAAAClgKgAAACAFTAUAAACkgKkAAACAFDAVAAAAkAKmAgAAAFLAVAAAAEAKmAoAAABIAVMBAAAAKWAqAAAAIAVMBQAAAKSAqQAAAIAUMBUAAACQAqYCAAAAUsBUAAAAQAqYCgAAAEgBUwEAAAApYCoAAAAgBUwFAAAApICpAAAAgBQwFQAAAJACpgIAAABSwFQAAABACpgKAAAASAFTAQAAAClgKgAAACAFTAUAAACkgKkAAACAFDAVAAAAkAKmAgAAAFLAVAAAAEAKmAoAAABIAVMBAAAAKWAqAAAAIAVMBQAAAKSAqQAAAIAU/hXB5Bfnt0ftGAAAAABJRU5ErkJggg=="


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("input")
    parser.add_argument("output")
    args = parser.parse_args()
    out = render(args.input, args.output)
    print(f"Wrote {out} ({out.stat().st_size:,} bytes)")
